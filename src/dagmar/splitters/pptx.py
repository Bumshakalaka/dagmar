"""PowerPoint splitter using LLM vision model for image description.

This module provides functionality to convert PowerPoint presentations to markdown
using LLM vision models to describe images inline, preserving slide structure.
"""

import base64
import io
import logging
import os
import re
from dataclasses import dataclass
from pathlib import Path
from typing import List

from dotenv import find_dotenv, load_dotenv
from langchain.chat_models import init_chat_model
from langchain_core.documents import Document
from langchain_text_splitters import MarkdownHeaderTextSplitter
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from dagmar.md_fixer import MarkdownFixer
from dagmar.splitters.base import BaseSplitter

logger = logging.getLogger(__name__)

# Load environment variables
load_dotenv(find_dotenv())


@dataclass(eq=False)
class PptxLlmSplitter(BaseSplitter):
    """Splits PowerPoint files using LLM vision model for image descriptions.

    This class converts PowerPoint presentations to markdown, using LLM vision
    models to describe images inline rather than extracting them as files.
    Each slide becomes one page in the final document.
    """

    file_pattern_re = r".+\.pptx"
    priority: int = 10
    use_cache: bool = True

    # LLM prompt for image description
    EXTRACTION_PROMPT = (Path(__file__).parent / "image_to_md_prompt.md").read_text()

    @classmethod
    def _extract_text_from_shape(cls, shape) -> str:
        """Extract text content from a shape."""
        if not hasattr(shape, "text_frame"):
            return ""

        text_frame = getattr(shape, "text_frame", None)
        if not text_frame:
            return ""

        text_parts = []
        for paragraph in text_frame.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text.strip())

        return "\n".join(text_parts)

    @classmethod
    def _normalize_rotation(cls, rotation_angle: float) -> float:
        """Normalize rotation angle to 0, 90, 180, or 270 degrees."""
        if rotation_angle is None:
            return 0.0

        # Normalize to 0-360 range
        normalized = rotation_angle % 360

        # Round to nearest 90-degree increment
        if normalized <= 45:
            return 0.0
        elif normalized <= 135:
            return 90.0
        elif normalized <= 225:
            return 180.0
        elif normalized <= 315:
            return 270.0
        else:
            return 0.0

    @classmethod
    def _extract_table_content(cls, table) -> str:
        """Extract table content and format as Markdown table."""
        if not table.rows:
            return ""

        markdown_table = []

        # Extract header row
        header_row = []
        if table.rows:
            first_row = table.rows[0]
            for cell in first_row.cells:
                cell_text = cell.text.strip() if cell.text else ""
                header_row.append(cell_text.replace("\n", "<br>"))

            if any(header_row):  # Only add header if it contains content
                markdown_table.append("| " + " | ".join(header_row) + " |")
                markdown_table.append("| " + " | ".join(["---"] * len(header_row)) + " |")

        # Extract data rows
        start_row = 1 if any(header_row) else 0
        for row_idx in range(start_row, len(table.rows)):
            row = table.rows[row_idx]
            row_data = []
            for cell in row.cells:
                cell_text = cell.text.strip() if cell.text else ""
                cell_text = cell_text.replace("\n", "<br>")
                row_data.append(cell_text)

            if any(row_data):  # Only add row if it contains content
                markdown_table.append("| " + " | ".join(row_data) + " |")

        return "\n".join(markdown_table) if markdown_table else ""

    @classmethod
    def _resize_image(cls, image: Image.Image, max_dimension: int = 1024) -> Image.Image:
        """Resize image to fit within max_dimension while preserving aspect ratio.

        Only downsamples images; does not upscale.

        :param image: PIL Image object to resize.
        :param max_dimension: Maximum allowed dimension (width or height). Default 1024px.
        :return: Resized PIL Image object.
        """
        width, height = image.size

        # Check if resizing is needed (only downsize)
        if width <= max_dimension and height <= max_dimension:
            return image

        # Calculate scaling factor to fit within max_dimension
        scale = min(max_dimension / width, max_dimension / height)
        new_width = int(width * scale)
        new_height = int(height * scale)

        logger.debug(f"Resizing image from {width}x{height} to {new_width}x{new_height}")
        return image.resize((new_width, new_height), Image.Resampling.LANCZOS)

    @classmethod
    def _process_image_with_llm(cls, image_blob: bytes, slide_num: int, img_num: int, chat_model) -> str:
        """Process a single image with LLM vision model.

        :param image_blob: Raw image bytes.
        :param slide_num: Slide number (1-indexed).
        :param img_num: Image number on slide (1-indexed).
        :param chat_model: Initialized chat model with vision capabilities.
        :return: Extracted markdown content for the image.
        """
        logger.debug(f"Processing image {img_num} on slide {slide_num} with LLM")
        try:
            # Load image from blob
            with Image.open(io.BytesIO(image_blob)) as img:
                # Resize if needed
                resized_image = cls._resize_image(img)

                # Convert to JPEG for LLM
                buffered = io.BytesIO()
                if resized_image.mode in ("RGBA", "LA"):
                    # Create white background for transparent images
                    background = Image.new("RGB", resized_image.size, (255, 255, 255))
                    if resized_image.mode == "RGBA":
                        background.paste(resized_image, mask=resized_image.split()[-1])
                    else:
                        background.paste(resized_image)
                    background.save(buffered, format="JPEG")
                else:
                    resized_image.save(buffered, format="JPEG")

                image_base64 = base64.b64encode(buffered.getvalue()).decode()

            # Create message with image
            messages = [
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": cls.EXTRACTION_PROMPT},
                        {
                            "type": "image_url",
                            "image_url": {"url": f"data:image/jpeg;base64,{image_base64}"},
                        },
                    ],
                }
            ]

            # Invoke the model
            response = chat_model.invoke(messages)
            markdown_content = response.content
            logger.debug(f"Successfully processed image {img_num} on slide {slide_num}")

            return markdown_content

        except Exception as e:
            logger.error(f"Failed to process image {img_num} on slide {slide_num}: {e}")
            return f"[ERROR: Image {img_num} on slide {slide_num} failed - {str(e)}]"

    @classmethod
    def _get_shape_position(cls, shape) -> tuple:
        """Get shape position for sorting (top, then left)."""
        top = shape.top / Inches(1) if hasattr(shape, "top") and shape.top is not None else 0.0
        left = shape.left / Inches(1) if hasattr(shape, "left") and shape.left is not None else 0.0
        return (top, left)

    @classmethod
    def _process_slide(cls, slide, slide_num: int, chat_model) -> str:
        """Process a single slide and return markdown content.

        :param slide: PowerPoint slide object.
        :param slide_num: Slide number (1-indexed).
        :param chat_model: Initialized chat model with vision capabilities.
        :return: Markdown content for the slide.
        """
        logger.debug(f"Processing slide {slide_num}")

        # Collect shapes with their positions for sorting
        shapes_with_positions = []
        for shape in slide.shapes:
            position = cls._get_shape_position(shape)
            shapes_with_positions.append((position, shape))

        # Sort by position (top first, then left)
        shapes_with_positions.sort(key=lambda x: x[0])

        # Process shapes in reading order
        slide_text_content = []
        slide_tables = []
        slide_images = []
        image_counter = 0

        for _, shape in shapes_with_positions:
            try:
                # Handle text content
                if hasattr(shape, "text_frame") and getattr(shape, "text_frame", None):
                    text_content = cls._extract_text_from_shape(shape)
                    if text_content:
                        slide_text_content.append(text_content)

                # Handle tables
                if hasattr(shape, "table") and getattr(shape, "table", None):
                    table_content = cls._extract_table_content(getattr(shape, "table"))
                    if table_content:
                        slide_tables.append(table_content)

                # Handle images
                if hasattr(shape, "image") and getattr(shape, "image", None):
                    image_counter += 1
                    image_blob = shape.image.blob

                    # Process image with LLM
                    image_markdown = cls._process_image_with_llm(image_blob, slide_num, image_counter, chat_model)
                    slide_images.append(image_markdown)

            except Exception as e:
                logger.warning(f"Error processing shape on slide {slide_num}: {e}")
                continue

        # Build slide markdown
        slide_markdown = []

        # Determine slide title
        slide_title = f"## Slide {slide_num}"
        if (
            slide_text_content
            and isinstance(slide_text_content[0], str)
            and len(slide_text_content[0]) < 100
            and "\n" not in slide_text_content[0]
        ):
            slide_title = f"## {slide_text_content[0]}"
            slide_text_content = slide_text_content[1:]  # Remove title from content

        slide_markdown.append(slide_title)
        slide_markdown.append("")

        # Add text content
        for text in slide_text_content:
            slide_markdown.append(text)
            slide_markdown.append("")

        # Add tables
        for table in slide_tables:
            slide_markdown.append(table)
            slide_markdown.append("")

        # Add images (LLM-described content)
        for image_md in slide_images:
            slide_markdown.append(image_md)
            slide_markdown.append("")

        return "\n".join(slide_markdown)

    @classmethod
    def split(cls, file_path: str) -> List[Document]:
        """Split a PowerPoint file into documents using LLM-based vision extraction.

        Converts PowerPoint to markdown with LLM-described images, then splits
        the result into documents. Uses caching to avoid reprocessing unchanged files.

        :param file_path: Path to the PowerPoint file to be split.
        :return: A list of Document objects resulting from the split.
        """
        logger.info(f"Processing PowerPoint file with LLM vision: {file_path}")
        try:
            # Check cache first if caching is enabled
            combined_markdown = None
            if cls.use_cache:
                combined_markdown = cls._load_from_cache(file_path)
                if combined_markdown is not None:
                    logger.info(f"Using cached content for {file_path}")
                else:
                    logger.debug(f"No cache found for {file_path}")

            # If no cache hit, process with LLM
            if combined_markdown is None:
                # Load presentation
                logger.debug("Loading PowerPoint presentation")
                presentation = Presentation(file_path)
                logger.debug(f"Loaded {len(presentation.slides)} slides")

                # Initialize LLM model
                if os.getenv("AZURE_OPENAI_ENDPOINT"):
                    logger.debug("Using Azure OpenAI model")
                    chat_model = init_chat_model(
                        model="gpt-4.1",
                        model_provider="azure_openai",
                        azure_endpoint=os.getenv("AZURE_OPENAI_ENDPOINT"),
                        api_key=os.getenv("AZURE_OPENAI_API_KEY"),
                        api_version=os.getenv("OPENAI_API_VERSION"),
                        max_tokens=16384,
                    )
                else:
                    logger.debug("Using OpenAI model")
                    chat_model = init_chat_model(
                        model="gpt-4.1",
                        model_provider="openai",
                        api_key=os.getenv("OPENAI_API_KEY"),
                    )

                # Process all slides
                logger.debug("Processing slides with LLM vision model")
                slide_markdowns = []
                for slide_idx, slide in enumerate(presentation.slides, 1):
                    # Add page marker
                    slide_md = cls._process_slide(slide, slide_idx, chat_model)
                    slide_md += f"\n<!-- PAGE:{slide_idx} -->"
                    slide_markdowns.append(slide_md)

                combined_markdown = "\n\n".join(slide_markdowns)

                # Save to cache if caching is enabled
                if cls.use_cache:
                    cls._save_to_cache(file_path, combined_markdown)

            # Fix markdown
            logger.debug("Fixing markdown")
            fixer = MarkdownFixer()
            combined_markdown = fixer.process_content(combined_markdown, Path(file_path).stem)

            # Split by H2 headers to preserve semantic sections
            text_splitter = MarkdownHeaderTextSplitter(headers_to_split_on=[("##", "H2")], strip_headers=False)
            documents = text_splitter.split_text(combined_markdown)

            # Extract page markers per chunk, add metadata, and strip markers from content
            prev_page_ids = set()
            for doc in documents:
                doc.metadata["source"] = file_path

                page_ids = set()
                for m in re.findall(r"<!--\s*PAGE:(\d+)\s*-->", doc.page_content):
                    page_ids.add(int(m))

                if page_ids:
                    pages_sorted = sorted(page_ids)
                    prev_page_ids = pages_sorted
                    doc.metadata["page_nums"] = pages_sorted
                elif prev_page_ids:
                    doc.metadata["page_nums"] = prev_page_ids

                # Remove markers so they don't affect embeddings
                doc.page_content = re.sub(r"<!--\s*PAGE:\d+\s*-->\s*", "", doc.page_content)

            logger.debug(f"PowerPoint file split into {len(documents)} chunks")
            return documents

        except Exception as e:
            logger.error(f"Failed to process PowerPoint file {file_path}: {e}")
            raise
